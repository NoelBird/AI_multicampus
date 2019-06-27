import sys
import os
import markdown as md
from pptx import Presentation
from bs4 import BeautifulSoup

from PyQt5 import QtWidgets
from PyQt5.QtGui import QIcon, QPixmap, QClipboard, QImage
from PyQt5.QtWidgets import QApplication, QWidget, QInputDialog, QLineEdit, QFileDialog, QPushButton, QHBoxLayout, QDialog, QBoxLayout, QMainWindow, QLabel
from PyQt5.QtWebEngineWidgets import QWebEnginePage
from PyQt5.QtWebEngineWidgets import QWebEngineView
from PyQt5.QtCore import QUrl, QDir, QMimeData
from formpy.ui_mainwindow2 import Ui_mainWindow
import PyQt5.QtCore as QtCore
from formpy.ui_full_disp import Ui_Dialog


class Mainwin(QMainWindow, Ui_mainWindow):
    def __init__(self):
        super(Mainwin, self).__init__()
        QMainWindow.__init__(self, parent=None)
        self.setupUi(self)
        self.fname = ''
        self.template = 'github.css'
        self.keylist = []

    def keyPressEvent(self, event): # clipboard to image file
        if event.key() == QtCore.Qt.Key_F8:
            self.imgFromClipboard()



    def openFile(self):
        options = QFileDialog.Options()
        # options |= QFileDialog.DontUseNativeDialog
        # TODO: 상속하는 부분 개선하기 None -> ?
        fileName, _ = QFileDialog.getOpenFileName(None, "열기", "","Markdown(*.md);;Power Point(*.pptx | *.ppt);;All Files (*)", options=options)

        if fileName:
            # TODO: file에서 dot(.)이 여러개일 경우에 에러가 없는지 체크
            ext = os.path.splitext(fileName)[1]
            if ext == '.md':
                self.loadMD(fileName)
            # if ext.lower() in ['.ppt', '.pptx']:

    # TODO: 특정한 디렉토리 내에서는 오픈이 안되는 이슈가 있습니다.
    def loadMD(self, fname):
        print(fname)
        with open(fname, "r") as fp:
            txt = fp.read()
            self.txtMarkdown.setText(txt)

    def saveFile(self):
        options = QFileDialog.Options()
        # options |= QFileDialog.DontUseNativeDialog
        # TODO: 상속하는 부분 개선하기 None -> ?
        fileName, _ = QFileDialog.getSaveFileName(None, "저장", "",
                                                  "Markdown(*.md);;Power Point(*.pptx | *.ppt);;All Files (*)",
                                                  options=options)

        if fileName:
            # TODO: file에서 dot(.)이 여러개일 경우에 에러가 없는지 체크
            ext = os.path.splitext(fileName)[1]
            if ext == '.md':
                self.saveMD(fileName)
            # if ext.lower() in ['.ppt', '.pptx']:

    def saveMD(self, fname):
        with open(fname, "wt") as fp:
            fp.write(self.txtMarkdown.toPlainText())

    def exportPPT(self):
        options = QFileDialog.Options()
        # options |= QFileDialog.DontUseNativeDialog
        # TODO: 상속하는 부분 개선하기 None -> ?
        fileName, _ = QFileDialog.getSaveFileName(None, "저장", "",
                                                  "Power Point(*.pptx | *.ppt);;All Files (*)",
                                                  options=options)
        htmlLst = md.markdown(self.txtMarkdown.toPlainText()).split('<hr />')
        prs = Presentation()
        for i in range(len(htmlLst)):
            bs = BeautifulSoup(htmlLst[i], "html.parser")
            title_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            # title 부분

            title.text = bs.find('h1').get_text()
            # subtitle.text = pTitle.findall(txtLst[i])[0]

            # subtitle 부분

            subtitle.text = bs.find('h2').get_text()

            # 본문 부분
            shapes = slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            p = tf.add_paragraph()
            # tf.text = '\n'.join([k for k in bs.findAll('p')])
            p.text = '안녕하세요'

        prs.save(fileName)

    def importPPT(self):
        pass

    def setTemplateGithub(self):
        self.template = 'github.css'
        self.updateWeb()

    def setTemplateDracula(self):
        self.template = 'dracula.css'
        self.updateWeb()

    def setTemplatePresentation(self):
        self.template = 'presentation.css'
        self.updateWeb()

    def setTemplateTmp(self):
        self.template = 'tmp.css'
        self.updateWeb()

    def imgFromClipboard(self):
        clipboard = QApplication.clipboard()
        mimedata = clipboard.mimeData()
        if(mimedata.hasImage()):
            lblImage = QLabel()
            img = QImage(mimedata.imageData())
            img.save('fff.png')
            self.txtMarkdown.append("![](fff.png)")
            # lblImage.setPixmap(mimedata.imageData())




    # def savePDF(self, fname):
    #     pdfkit.from_file(self.updatePreview(), 'out.pdf')

    def close(self):
        sys.exit()


    def displayFullScreen(self):
        import copy
        wgt = QDialog()
        wgt.showFullScreen()
        layout = QBoxLayout(QBoxLayout.TopToBottom)
        wgt.setLayout(layout)
        layout.addWidget(self.wgtWeb) # widget이 새 창꺼랑 같이 꺼지는 경향이 있음
        # wgt.setWindowState()
        wgt.exec_()



    # def keyPressEvent(self, event):
    #     if event.key() == QtCore.Qt.Key_Control | QtCore.Qt.Key_V:
    #         print("here")
    #     event.accept()

    def updateWeb(self):
        try:
            curMD = md.markdown(self.txtMarkdown.toPlainText())
            baseDir = "C:\\Users\\user\\Documents\\Python-programming\\quick-point"
            cssFile = self.template
            css = "<link rel=\"stylesheet\" type=\"text/css\" href=\"templates/%s\"><br>" % cssFile
            totalHtml = "<div class='markdown-body'>" + css + '\n' + curMD +"</div>"
            with open("curHtml.html", "wt") as fp:
                fp.write(totalHtml)
            print(css+curMD)
            # self.wgtWeb.setHtml(css+curMD)
            # print(QDir.currentPath() + '/curHtml.html')
            self.wgtWeb.setUrl(QUrl.fromLocalFile(QDir.currentPath() + '/curHtml.html'))
        except Exception as e:
            pass
        finally:
            pass
            # pixmap = QPixmap('out.jpg')
            # self.lblPreview.setFixedSize(self.lblPreview.size())
            # self.lblPreview.setPixmap(pixmap)

#
# class ctrlWin(QtWidgets.QMainWindow):
#     cnt = 0
#     def __init__(self):
#         self.win = QtWidgets.QMainWindow()
#         ui = Mainwin()
#         ui.setupUi(self.win)
#         self.win.setWindowIcon(QIcon('logoico.ico'))
#         self.win.statusBar().showMessage('준비')
#         self.wgtWeb = self.win.findChild(QWebEngineView)
#
#         self.win.txtMarkdown.connect()
#         self.win.show()
#
#     def updateWindow(self):
#         self.wgtWeb.setHtml("%d" % ctrlWin.cnt)
#         ctrlWin += 1
#         # submitFrame = self.subwindow.findChild(QtGui.QFrame, "frameSubmit")


if __name__ == '__main__':
    import sys
    app = QtWidgets.QApplication(sys.argv)
    MainWindow = Mainwin()
    MainWindow.setWindowIcon(QIcon('logoico.ico'))
    MainWindow.statusBar().showMessage('준비')
    MainWindow.show()
    sys.exit(app.exec_())
